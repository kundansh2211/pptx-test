import os
from git import Repo
from pptx import Presentation
from pptx.util import Inches

def get_images_from_last_n_commits(repo_path, n):
    # Initialize the repository
    repo = Repo(repo_path)

    # Get the list of the last n commits
    commits = list(repo.iter_commits('main', max_count=n))  # Replace 'main' with the correct branch if necessary

    image_files = set()

    # Loop through each commit
    for commit in commits:
        # Get the diff of the commit
        for diff in commit.diff(None, create_patch=False):
            # Check if the file is an image
            if diff.a_path and diff.a_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                image_files.add(diff.a_path)
            if diff.b_path and diff.b_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                image_files.add(diff.b_path)

    return list(image_files)

def create_presentation(image_list, output_path):
    # Create a presentation object
    presentation = Presentation()

    # Loop through the image filenames and add slides with images
    for image_path in image_list:
        # Check if the image file exists
        if not os.path.isfile(image_path):
            print(f"Image file not found, skipping: {image_path}")
            continue

        try:
            # Create a slide with blank layout
            slide_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(slide_layout)

            # Add the image to the slide and adjust size and position
            pic = slide.shapes.add_picture(image_path, 0, 0)

            # Calculate the scaling factor to maintain the aspect ratio
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height
            scale_w = slide_width / pic.width
            scale_h = slide_height / pic.height
            scale = min(scale_w, scale_h)

            # Resize the image
            new_width = int(pic.width * scale)
            new_height = int(pic.height * scale)

            # Center the image on the slide
            pic.left = int((slide_width - new_width) / 2)
            pic.top = int((slide_height - new_height) / 2)
            pic.width = new_width
            pic.height = new_height

        except Exception as e:
            print(f"Error adding image {image_path}: {e}")

    # Save the presentation
    presentation.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

    # Cleanup temporary files
    cleanup_temporary_files(os.path.dirname(output_path))

def cleanup_temporary_files(directory):
    for filename in os.listdir(directory):
        if filename.startswith('~$'):
            tmp_file_path = os.path.join(directory, filename)
            try:
                os.remove(tmp_file_path)
                print(f"Removed temporary file: {tmp_file_path}")
            except Exception as e:
                print(f"Error removing temporary file {tmp_file_path}: {e}")

# Example usage:
repo_path = '.'  # Current directory, assuming the script is in the repo
n = 5  # Number of commits to look back
image_list = get_images_from_last_n_commits(repo_path, n)
print("Images from last", n, "commits:", image_list)

# Ensure the output directory exists
output_dir = 'PPT'
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Set the output path for the presentation
output_path = os.path.join(output_dir, 'image_ppt.pptx')

# Create the presentation
create_presentation(image_list, output_path)
