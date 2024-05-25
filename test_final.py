import os
from pptx import Presentation
from pptx.util import Inches

# Define the directory where images are stored
image_dir = 'images'

# List of images to be converted to ppt
image_filenames = ['image1.jpg', 'image2.jpg']

# Check if the image directory exists
if not os.path.exists(image_dir):
    print(f"Image directory not found: {image_dir}")
    exit()

# Create a presentation object
presentation = Presentation()

# Loop through the image filenames and add slides with images
for filename in image_filenames:
    # Construct the full path to the image file
    image_path = os.path.join(image_dir, filename)

    # Check if the file exists
    if not os.path.isfile(image_path):
        print(f"File not found: {image_path}")
        continue

    try:
        # Create a slide with blank layout
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        # Add the image to the slide and adjust size and position
        pic = slide.shapes.add_picture(image_path, 0, 0)
        
        # Calculate the scaling factor to maintain the aspect ratio
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        scale_w = slide_width / pic.width
        scale_h = slide_height / pic.height
        scale = min(scale_w, scale_h)

        # Resize the image
        new_width = int(pic.width * scale)
        new_height = int(pic.height * scale)

        # Center the image on the slide
        pic.left = int((slide_width - new_width) / 2)
        pic.top = int((slide_height - new_height) / 2)
        pic.width = new_width
        pic.height = new_height

    except Exception as e:
        print(f"Error adding image {image_path}: {e}")

# Ensure the output directory exists
output_dir = 'PPT'
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Save the presentation
output_path = os.path.join(output_dir, 'image_ppt.pptx')
presentation.save(output_path)
print(f"Presentation saved successfully to {output_path}")
